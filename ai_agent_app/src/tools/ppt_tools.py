"""
PowerPoint Automation Tools
Uses python-pptx for all PowerPoint operations.
"""
import os
from contextlib import contextmanager
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from datetime import datetime

try:
    import pythoncom
    from win32com.client import DispatchEx, GetActiveObject
except Exception:
    pythoncom = None
    DispatchEx = None
    GetActiveObject = None


def _hex_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _normalize_hex_color(hex_color: str, default: str) -> str:
    cleaned = (hex_color or "").strip().lstrip("#")
    if len(cleaned) == 6:
        return cleaned.upper()
    return default


def _open_created_presentation(filename: str) -> str:
    """Open a saved PowerPoint file with the default Windows app when possible."""
    try:
        absolute = os.path.abspath(filename)
        os.startfile(absolute)
        return f" Opened in PowerPoint: {absolute}"
    except Exception:
        return ""


def _normalize_presentation_path(filename: str) -> str:
    """Ensure PowerPoint files use a real .pptx path, even for __ACTIVE__ placeholder inputs."""
    cleaned = (filename or "").strip()
    if not cleaned or cleaned.upper() == "__ACTIVE__":
        cleaned = "__ACTIVE__.pptx"
    if not cleaned.lower().endswith(".pptx"):
        cleaned = f"{cleaned}.pptx"
    return cleaned


def _is_active_presentation_target(filename: str) -> bool:
    cleaned = (filename or "").strip().upper()
    return cleaned in {"", "__ACTIVE__", "ACTIVE", "__ACTIVE__.PPTX"}


@contextmanager
def _com_context():
    if pythoncom is None:
        yield False
        return
    pythoncom.CoInitialize()
    try:
        yield True
    finally:
        pythoncom.CoUninitialize()


def _get_powerpoint_app(visible: bool = True):
    if pythoncom is None or DispatchEx is None or GetActiveObject is None:
        raise RuntimeError("PowerPoint desktop automation requires Windows with pywin32 and Microsoft PowerPoint installed.")
    try:
        app = GetActiveObject("PowerPoint.Application")
    except Exception:
        app = DispatchEx("PowerPoint.Application")
    if visible:
        app.Visible = True
    return app


def _find_open_presentation(app, filename: str):
    target = os.path.abspath(filename).lower()
    for presentation in app.Presentations:
        try:
            if os.path.abspath(str(presentation.FullName)).lower() == target:
                return presentation
        except Exception:
            continue
    return None


def _get_powerpoint_presentation(app, filename: str, create_if_missing: bool = False):
    if _is_active_presentation_target(filename):
        try:
            presentation = app.ActivePresentation
            if presentation:
                return presentation
        except Exception:
            pass
        if not create_if_missing:
            raise RuntimeError("No active PowerPoint presentation is open.")
        filename = _normalize_presentation_path(filename)
    else:
        filename = _normalize_presentation_path(filename)

    existing = _find_open_presentation(app, filename)
    if existing is not None:
        return existing

    absolute = os.path.abspath(filename)
    if os.path.exists(absolute):
        return app.Presentations.Open(absolute, WithWindow=True)

    if not create_if_missing:
        raise FileNotFoundError(f"PowerPoint file not found: {absolute}")

    os.makedirs(os.path.dirname(absolute), exist_ok=True)
    presentation = app.Presentations.Add(WithWindow=True)
    presentation.SaveAs(absolute)
    return presentation


def _apply_com_font(text_range, *, font_size: int = None, font_color: str = "", font_name: str = "", bold: bool = None):
    font = text_range.Font
    if font_size:
        font.Size = max(font_size, 1)
    if font_color:
        rgb = _normalize_hex_color(font_color, "")
        if rgb:
            font.Color.RGB = int(rgb[4:6] + rgb[2:4] + rgb[0:2], 16)
    if font_name:
        font.Name = font_name
    if bold is not None:
        font.Bold = -1 if bold else 0


def _add_slide_via_com(filename: str, title: str, content: list = None,
                       title_font_size: int = 28, content_font_size: int = 18,
                       title_font_color: str = "FFFFFF", content_font_color: str = "0F172A",
                       background_color: str = "F1F5F9", fill_color: str = "",
                       title_font_name: str = "", content_font_name: str = "") -> str:
    with _com_context():
        app = _get_powerpoint_app(visible=True)
        presentation = _get_powerpoint_presentation(app, filename, create_if_missing=False)
        slide = presentation.Slides.Add(presentation.Slides.Count + 1, 12)
        slide.Background.Fill.Solid()
        slide.Background.Fill.ForeColor.RGB = int(_normalize_hex_color(fill_color or background_color, "F1F5F9")[4:6] + _normalize_hex_color(fill_color or background_color, "F1F5F9")[2:4] + _normalize_hex_color(fill_color or background_color, "F1F5F9")[0:2], 16)

        title_box = slide.Shapes.AddTextbox(1, 28, 14, 680, 40)
        title_range = title_box.TextFrame.TextRange
        title_range.Text = title
        _apply_com_font(title_range, font_size=title_font_size, font_color=title_font_color, font_name=title_font_name, bold=True)

        if content:
            body_box = slide.Shapes.AddTextbox(1, 36, 94, 660, 280)
            body_range = body_box.TextFrame.TextRange
            body_range.Text = "\r\n".join(f"- {item}" for item in content)
            _apply_com_font(body_range, font_size=content_font_size, font_color=content_font_color, font_name=content_font_name)

        presentation.Save()
        return f"✅ Slide '{title}' added to {presentation.FullName} (total slides: {presentation.Slides.Count})"


def _edit_slide_via_com(filename: str, slide_index: int, new_title: str = None, new_content: list = None,
                        title_font_size: int = 0, content_font_size: int = 0,
                        title_font_color: str = "", content_font_color: str = "",
                        title_font_name: str = "", content_font_name: str = "") -> str:
    with _com_context():
        app = _get_powerpoint_app(visible=True)
        presentation = _get_powerpoint_presentation(app, filename, create_if_missing=False)
        if slide_index < 0 or slide_index >= presentation.Slides.Count:
            return f"❌ Slide index {slide_index} out of range. File has {presentation.Slides.Count} slides."

        slide = presentation.Slides(slide_index + 1)
        text_shapes = []
        for shape in slide.Shapes:
            try:
                if shape.HasTextFrame:
                    text_shapes.append(shape)
            except Exception:
                continue

        title_shape = text_shapes[0] if text_shapes else None
        body_shapes = text_shapes[1:] if len(text_shapes) > 1 else []

        if title_shape is not None:
            title_range = title_shape.TextFrame.TextRange
            if new_title:
                title_range.Text = new_title
            _apply_com_font(
                title_range,
                font_size=title_font_size or None,
                font_color=title_font_color,
                font_name=title_font_name,
                bold=True,
            )

        if new_content is not None:
            body_text = "\r\n".join(f"- {item}" for item in new_content)
            if body_shapes:
                body_range = body_shapes[0].TextFrame.TextRange
            else:
                body_box = slide.Shapes.AddTextbox(1, 36, 94, 660, 280)
                body_range = body_box.TextFrame.TextRange
            body_range.Text = body_text
            _apply_com_font(
                body_range,
                font_size=content_font_size or None,
                font_color=content_font_color,
                font_name=content_font_name,
            )
        elif body_shapes:
            body_range = body_shapes[0].TextFrame.TextRange
            _apply_com_font(
                body_range,
                font_size=content_font_size or None,
                font_color=content_font_color,
                font_name=content_font_name,
            )

        presentation.Save()
        return f"✅ Slide {slide_index} edited in {presentation.FullName}"


def _add_table_slide_via_com(filename: str, title: str, headers: list, rows: list = None,
                             title_font_size: int = 28, cell_font_size: int = 16,
                             title_font_color: str = "FFFFFF", cell_font_color: str = "0F172A",
                             background_color: str = "F1F5F9", fill_color: str = "",
                             title_font_name: str = "", cell_font_name: str = "") -> str:
    rows = rows or []
    with _com_context():
        app = _get_powerpoint_app(visible=True)
        presentation = _get_powerpoint_presentation(app, filename, create_if_missing=False)
        slide = presentation.Slides.Add(presentation.Slides.Count + 1, 12)
        slide.Background.Fill.Solid()
        bg_hex = _normalize_hex_color(fill_color or background_color, "F1F5F9")
        slide.Background.Fill.ForeColor.RGB = int(bg_hex[4:6] + bg_hex[2:4] + bg_hex[0:2], 16)

        title_box = slide.Shapes.AddTextbox(1, 28, 14, 680, 40)
        title_range = title_box.TextFrame.TextRange
        title_range.Text = title
        _apply_com_font(title_range, font_size=title_font_size, font_color=title_font_color, font_name=title_font_name, bold=True)

        total_rows = 1 + len(rows)
        total_cols = len(headers)
        table_shape = slide.Shapes.AddTable(total_rows, total_cols, 30, 95, 650, 260)
        table = table_shape.Table

        for col_idx, header in enumerate(headers, start=1):
            cell_range = table.Cell(1, col_idx).Shape.TextFrame.TextRange
            cell_range.Text = str(header)
            _apply_com_font(cell_range, font_size=cell_font_size, font_color=cell_font_color, font_name=cell_font_name, bold=True)

        for row_idx, row_data in enumerate(rows, start=2):
            padded = list(row_data) + [""] * (total_cols - len(row_data))
            for col_idx, value in enumerate(padded[:total_cols], start=1):
                cell_range = table.Cell(row_idx, col_idx).Shape.TextFrame.TextRange
                cell_range.Text = str(value)
                _apply_com_font(cell_range, font_size=cell_font_size, font_color=cell_font_color, font_name=cell_font_name)

        presentation.Save()
        return f"✅ Table slide '{title}' added to {presentation.FullName} (total slides: {presentation.Slides.Count})"


def ppt_create_presentation(filename: str, title: str, subtitle: str = "",
                              company: str = "", date: str = "",
                              title_font_size: int = 44, subtitle_font_size: int = 22,
                              title_font_color: str = "FFFFFF", subtitle_font_color: str = "00B4D8",
                              background_color: str = "1A2B5E", fill_color: str = "",
                              title_font_name: str = "", subtitle_font_name: str = "") -> str:
    """Create a new PowerPoint with a professional title slide."""
    try:
        filename = _normalize_presentation_path(filename)
        os.makedirs(os.path.dirname(os.path.abspath(filename)), exist_ok=True)
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        slide_fill = _normalize_hex_color(fill_color or background_color, "1A2B5E")
        title_color = _normalize_hex_color(title_font_color, "FFFFFF")
        subtitle_color = _normalize_hex_color(subtitle_font_color, "00B4D8")

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        W = prs.slide_width
        H = prs.slide_height

        # Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex_rgb(slide_fill)

        # Left accent bar
        from pptx.util import Emu
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), H)
        bar.fill.solid()
        bar.fill.fore_color.rgb = _hex_rgb("00B4D8")
        bar.line.fill.background()

        # Title
        txb = slide.shapes.add_textbox(Inches(0.4), Inches(1.5), Inches(9), Inches(1.2))
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(max(title_font_size, 1))
        run.font.bold = True
        run.font.color.rgb = _hex_rgb(title_color)
        if title_font_name:
            run.font.name = title_font_name

        # Subtitle
        if subtitle:
            txb2 = slide.shapes.add_textbox(Inches(0.4), Inches(2.85), Inches(9), Inches(0.7))
            tf2 = txb2.text_frame
            p2 = tf2.paragraphs[0]
            r2 = p2.add_run()
            r2.text = subtitle
            r2.font.size = Pt(max(subtitle_font_size, 1))
            r2.font.color.rgb = _hex_rgb(subtitle_color)
            if subtitle_font_name:
                r2.font.name = subtitle_font_name

        # Divider line
        from pptx.util import Pt as PtU
        line = slide.shapes.add_shape(1, Inches(0.4), Inches(3.7), Inches(6), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = _hex_rgb("1E6FBA")
        line.line.fill.background()

        # Company / Date row
        info_parts = []
        if company:
            info_parts.append(company)
        if date:
            info_parts.append(date)
        else:
            info_parts.append(datetime.now().strftime("%B %Y"))

        if info_parts:
            txb3 = slide.shapes.add_textbox(Inches(0.4), Inches(3.9), Inches(9), Inches(0.5))
            tf3 = txb3.text_frame
            p3 = tf3.paragraphs[0]
            r3 = p3.add_run()
            r3.text = "  ·  ".join(info_parts)
            r3.font.size = Pt(16)
            r3.font.color.rgb = _hex_rgb("CADCFC")

        prs.save(filename)
        opened_message = _open_created_presentation(filename)
        return f"✅ PowerPoint created: {filename}.{opened_message}"
    except Exception as e:
        return f"❌ Error creating PowerPoint: {e}"


def ppt_add_slide(filename: str, title: str, content: list = None,
                  layout: str = "title_content", title_font_size: int = 28,
                  content_font_size: int = 18, title_font_color: str = "FFFFFF",
                  content_font_color: str = "0F172A", background_color: str = "F1F5F9",
                  fill_color: str = "", title_font_name: str = "",
                  content_font_name: str = "") -> str:
    """Add a new content slide to an existing PowerPoint file."""
    try:
        filename = _normalize_presentation_path(filename)
        prs = Presentation(filename)
        W = prs.slide_width
        H = prs.slide_height
        slide_fill = _normalize_hex_color(fill_color or background_color, "F1F5F9")
        title_color = _normalize_hex_color(title_font_color, "FFFFFF")
        body_color = _normalize_hex_color(content_font_color, "0F172A")

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

        # Light background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex_rgb(slide_fill)

        # Header bar
        hdr = slide.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(1.1))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = _hex_rgb("1A2B5E")
        hdr.line.fill.background()

        # Title text
        txb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.75))
        tf = txb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(max(title_font_size, 1))
        run.font.bold = True
        run.font.color.rgb = _hex_rgb(title_color)
        if title_font_name:
            run.font.name = title_font_name

        # Content area
        if content:
            txb2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8))
            tf2 = txb2.text_frame
            tf2.word_wrap = True
            for i, item in enumerate(content):
                p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
                p2.space_after = Pt(8)
                run2 = p2.add_run()
                run2.text = f"•  {item}"
                run2.font.size = Pt(max(content_font_size, 1))
                run2.font.color.rgb = _hex_rgb(body_color)
                if content_font_name:
                    run2.font.name = content_font_name

        prs.save(filename)
        return f"✅ Slide '{title}' added to {filename} (total slides: {len(prs.slides)})"
    except PermissionError:
        try:
            return _add_slide_via_com(
                filename=filename,
                title=title,
                content=content,
                title_font_size=title_font_size,
                content_font_size=content_font_size,
                title_font_color=title_font_color,
                content_font_color=content_font_color,
                background_color=background_color,
                fill_color=fill_color,
                title_font_name=title_font_name,
                content_font_name=content_font_name,
            )
        except Exception as e:
            return f"❌ Error adding slide: {e}"
    except Exception as e:
        return f"❌ Error adding slide: {e}"


def ppt_add_table_slide(filename: str, title: str, headers: list,
                        rows: list = None, title_font_size: int = 28,
                        cell_font_size: int = 16, title_font_color: str = "FFFFFF",
                        cell_font_color: str = "0F172A", background_color: str = "F1F5F9",
                        fill_color: str = "", title_font_name: str = "",
                        cell_font_name: str = "") -> str:
    """Add a slide containing a native PowerPoint table."""
    try:
        filename = _normalize_presentation_path(filename)
        prs = Presentation(filename)
        rows = rows or []
        if not headers:
            return "❌ At least one table header is required."

        slide_fill = _normalize_hex_color(fill_color or background_color, "F1F5F9")
        title_color = _normalize_hex_color(title_font_color, "FFFFFF")
        cell_color = _normalize_hex_color(cell_font_color, "0F172A")

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex_rgb(slide_fill)

        header_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = _hex_rgb("1A2B5E")
        header_bar.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.75))
        title_frame = title_box.text_frame
        title_run = title_frame.paragraphs[0].add_run()
        title_run.text = title
        title_run.font.size = Pt(max(title_font_size, 1))
        title_run.font.bold = True
        title_run.font.color.rgb = _hex_rgb(title_color)
        if title_font_name:
            title_run.font.name = title_font_name

        total_rows = 1 + len(rows)
        total_cols = len(headers)
        table_shape = slide.shapes.add_table(total_rows, total_cols, Inches(0.45), Inches(1.45), Inches(12.2), Inches(4.8))
        table = table_shape.table

        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _hex_rgb("DCEBFA")
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(max(cell_font_size, 1))
                    run.font.bold = True
                    run.font.color.rgb = _hex_rgb(cell_color)
                    if cell_font_name:
                        run.font.name = cell_font_name

        for row_idx, row_data in enumerate(rows, start=1):
            padded = list(row_data) + [""] * (total_cols - len(row_data))
            for col_idx, value in enumerate(padded[:total_cols]):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)
                cell.fill.solid()
                cell.fill.fore_color.rgb = _hex_rgb("FFFFFF" if row_idx % 2 else "F8FAFC")
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(max(cell_font_size, 1))
                        run.font.color.rgb = _hex_rgb(cell_color)
                        if cell_font_name:
                            run.font.name = cell_font_name

        prs.save(filename)
        return f"✅ Table slide '{title}' added to {filename} (total slides: {len(prs.slides)})"
    except PermissionError:
        try:
            return _add_table_slide_via_com(
                filename=filename,
                title=title,
                headers=headers,
                rows=rows,
                title_font_size=title_font_size,
                cell_font_size=cell_font_size,
                title_font_color=title_font_color,
                cell_font_color=cell_font_color,
                background_color=background_color,
                fill_color=fill_color,
                title_font_name=title_font_name,
                cell_font_name=cell_font_name,
            )
        except Exception as e:
            return f"❌ Error adding table slide: {e}"
    except Exception as e:
        return f"❌ Error adding table slide: {e}"


def ppt_edit_slide(filename: str, slide_index: int,
                   new_title: str = None, new_content: list = None,
                   title_font_size: int = 0, content_font_size: int = 0,
                   title_font_color: str = "", content_font_color: str = "",
                   title_font_name: str = "", content_font_name: str = "") -> str:
    """Edit title/content of an existing slide."""
    try:
        filename = _normalize_presentation_path(filename)
        prs = Presentation(filename)
        if slide_index >= len(prs.slides):
            return f"❌ Slide index {slide_index} out of range. File has {len(prs.slides)} slides."

        slide = prs.slides[slide_index]
        text_shapes = [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False)]
        title_shape = text_shapes[0] if text_shapes else None
        body_shapes = text_shapes[1:] if len(text_shapes) > 1 else []

        if title_shape is not None:
            if new_title is not None:
                title_shape.text_frame.clear()
                para = title_shape.text_frame.paragraphs[0]
                run = para.add_run()
                run.text = new_title
            for para in title_shape.text_frame.paragraphs:
                for run in para.runs:
                    if title_font_size:
                        run.font.size = Pt(max(title_font_size, 1))
                    if title_font_color:
                        run.font.color.rgb = _hex_rgb(_normalize_hex_color(title_font_color, "000000"))
                    if title_font_name:
                        run.font.name = title_font_name
                    run.font.bold = True

        if new_content is not None:
            body_shape = body_shapes[0] if body_shapes else None
            if body_shape is None:
                body_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8))
            body_shape.text_frame.clear()
            for i, item in enumerate(new_content):
                para = body_shape.text_frame.paragraphs[0] if i == 0 else body_shape.text_frame.add_paragraph()
                run = para.add_run()
                run.text = f"•  {item}"
                if content_font_size:
                    run.font.size = Pt(max(content_font_size, 1))
                if content_font_color:
                    run.font.color.rgb = _hex_rgb(_normalize_hex_color(content_font_color, "0F172A"))
                if content_font_name:
                    run.font.name = content_font_name
        elif body_shapes:
            for para in body_shapes[0].text_frame.paragraphs:
                for run in para.runs:
                    if content_font_size:
                        run.font.size = Pt(max(content_font_size, 1))
                    if content_font_color:
                        run.font.color.rgb = _hex_rgb(_normalize_hex_color(content_font_color, "0F172A"))
                    if content_font_name:
                        run.font.name = content_font_name

        prs.save(filename)
        return f"✅ Slide {slide_index} edited in {filename}"
    except PermissionError:
        try:
            return _edit_slide_via_com(
                filename=filename,
                slide_index=slide_index,
                new_title=new_title,
                new_content=new_content,
                title_font_size=title_font_size,
                content_font_size=content_font_size,
                title_font_color=title_font_color,
                content_font_color=content_font_color,
                title_font_name=title_font_name,
                content_font_name=content_font_name,
            )
        except Exception as e:
            return f"❌ Error editing slide: {e}"
    except Exception as e:
        return f"❌ Error editing slide: {e}"
